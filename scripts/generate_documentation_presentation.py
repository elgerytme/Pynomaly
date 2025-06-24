#!/usr/bin/env python3
"""
Generate PowerPoint presentation from Pynomaly documentation.

This script creates a comprehensive PowerPoint presentation combining
all the documentation into a structured presentation suitable for
business stakeholders and technical audiences.
"""

import os
import sys
from pathlib import Path
from typing import List, Dict, Any
import json
import argparse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

class PynomlalyPresentationGenerator:
    """Generate PowerPoint presentations from Pynomaly documentation."""
    
    def __init__(self):
        self.presentation = Presentation()
        self.setup_styles()
        
    def setup_styles(self):
        """Setup presentation styles and themes."""
        # Set slide size to widescreen (16:9)
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)
        
        # Define color scheme
        self.colors = {
            'primary': RGBColor(0, 102, 204),      # Blue
            'secondary': RGBColor(51, 51, 51),     # Dark Gray
            'accent': RGBColor(255, 102, 0),       # Orange
            'success': RGBColor(40, 167, 69),      # Green
            'warning': RGBColor(255, 193, 7),      # Yellow
            'danger': RGBColor(220, 53, 69),       # Red
            'light': RGBColor(248, 249, 250),      # Light Gray
            'white': RGBColor(255, 255, 255)       # White
        }
    
    def create_title_slide(self):
        """Create the title slide."""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Pynomaly: State-of-the-Art Anomaly Detection Platform"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Comprehensive Platform Overview\nProcess, Architecture, Algorithms, and Business Procedures"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[1].font.size = Pt(16)
        subtitle.text_frame.paragraphs[1].font.color.rgb = self.colors['secondary']
        
        # Add footer with date and version
        self.add_footer(slide, "Pynomaly Documentation v1.0 - 2024")
    
    def create_agenda_slide(self):
        """Create agenda/overview slide."""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Presentation Agenda"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "1. Pynomaly Process Overview"
        
        agenda_items = [
            "2. System Architecture and Design",
            "3. Algorithm Options and Functionality", 
            "4. Autonomous Mode Capabilities",
            "5. Algorithm Selection Guidelines",
            "6. Business User Procedures",
            "7. Monthly Testing Workflows",
            "8. Implementation Roadmap",
            "9. Q&A and Discussion"
        ]
        
        for item in agenda_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
    
    def create_section_divider(self, section_title: str, section_subtitle: str = ""):
        """Create a section divider slide."""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.presentation.slide_width, 
            self.presentation.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['primary']
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.33), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.color.rgb = self.colors['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle if provided
        if section_subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(11.33), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = section_subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.colors['light']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def create_process_overview_slides(self):
        """Create slides for process overview section."""
        self.create_section_divider("Pynomaly Process Overview", "From Installation to Production")
        
        # Process workflow slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "End-to-End Process Workflow"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "1. Installation & Setup"
        
        process_steps = [
            "2. Data Preparation & Validation",
            "3. Algorithm Selection & Configuration",
            "4. Model Training & Optimization", 
            "5. Anomaly Detection & Analysis",
            "6. Results Interpretation & Reporting",
            "7. Model Evaluation & Validation",
            "8. Production Deployment",
            "9. Monitoring & Maintenance"
        ]
        
        for step in process_steps:
            p = text_frame.add_paragraph()
            p.text = step
            p.level = 0
            p.font.size = Pt(16)
        
        # Key benefits slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Key Process Benefits"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        benefits = [
            "• Streamlined workflow from data to insights",
            "• Automated quality checks and validations",
            "• Flexible deployment options (cloud, on-premise, edge)",
            "• Comprehensive monitoring and alerting",
            "• Integration with existing business systems",
            "• Scalable from prototype to production",
            "• Expert-level automation with user control",
            "• Detailed audit trails and compliance support"
        ]
        
        text_frame.text = benefits[0]
        for benefit in benefits[1:]:
            p = text_frame.add_paragraph()
            p.text = benefit
            p.font.size = Pt(16)
    
    def create_architecture_slides(self):
        """Create slides for architecture section."""
        self.create_section_divider("System Architecture", "Clean Architecture with Domain-Driven Design")
        
        # Architecture overview slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Layered Architecture Overview"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        layers = [
            "🎯 Presentation Layer: REST API, CLI, Web UI (PWA)",
            "⚙️ Application Layer: Use Cases, Services, DTOs", 
            "🏛️ Domain Layer: Entities, Value Objects, Business Logic",
            "🔧 Infrastructure Layer: Adapters, Persistence, External Services"
        ]
        
        text_frame.text = layers[0]
        for layer in layers[1:]:
            p = text_frame.add_paragraph()
            p.text = layer
            p.font.size = Pt(16)
        
        # Design principles slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Core Design Principles"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        principles = [
            "🎯 Separation of Concerns: Each layer has distinct responsibilities",
            "🔄 Dependency Inversion: High-level modules independent of low-level",
            "🧩 Interface Segregation: Clients depend only on what they use",
            "⚡ Single Responsibility: One reason to change per component",
            "🔓 Open/Closed Principle: Open for extension, closed for modification",
            "🏗️ Repository Pattern: Data access abstraction",
            "🏭 Factory Pattern: Object creation and algorithm instantiation",
            "🎭 Strategy Pattern: Interchangeable algorithms"
        ]
        
        text_frame.text = principles[0]
        for principle in principles[1:]:
            p = text_frame.add_paragraph()
            p.text = principle
            p.font.size = Pt(14)
    
    def create_algorithms_slides(self):
        """Create slides for algorithms section."""
        self.create_section_divider("Algorithm Options", "45+ Algorithms Across Multiple Categories")
        
        # Algorithm categories slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Algorithm Categories Overview"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        categories = [
            "📊 Statistical Methods (8): Isolation Forest, LOF, One-Class SVM",
            "🤖 Machine Learning (12): Random Forest, Gradient Boosting, k-NN",
            "🧠 Deep Learning (10): AutoEncoder, LSTM, CNN, Transformer",
            "🎯 Specialized (15): Graph Neural Networks, Time Series, Text Analysis",
            "🔗 Ensemble Methods (5): Voting, Stacking, Bagging, Adaptive"
        ]
        
        text_frame.text = categories[0]
        for category in categories[1:]:
            p = text_frame.add_paragraph()
            p.text = category
            p.font.size = Pt(16)
        
        # Performance comparison slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title  
        title.text = "Performance vs Resource Trade-offs"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        trade_offs = [
            "⚡ High Performance + Low Resource: Isolation Forest, Random Forest",
            "🎯 High Performance + High Resource: Deep Ensembles, Transformers",
            "⚖️ Medium Performance + Low Resource: Statistical Methods, PCA",
            "🔄 Medium Performance + Medium Resource: LOF, AutoEncoder"
        ]
        
        text_frame.text = trade_offs[0]
        for trade_off in trade_offs[1:]:
            p = text_frame.add_paragraph()
            p.text = trade_off
            p.font.size = Pt(16)
    
    def create_autonomous_mode_slides(self):
        """Create slides for autonomous mode section."""
        self.create_section_divider("Autonomous Mode", "AI-Powered Automated Anomaly Detection")
        
        # Autonomous capabilities slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Autonomous Mode Capabilities"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        capabilities = [
            "🔍 Automatic data analysis and characterization",
            "🎯 Intelligent algorithm selection and ranking",
            "⚙️ Automated hyperparameter optimization",
            "🧹 Smart data preprocessing and feature engineering",
            "📊 Continuous performance monitoring",
            "🔄 Adaptive learning and drift detection",
            "💡 Explainable AI decisions and recommendations",
            "🚀 Zero-configuration deployment option"
        ]
        
        text_frame.text = capabilities[0]
        for capability in capabilities[1:]:
            p = text_frame.add_paragraph()
            p.text = capability
            p.font.size = Pt(16)
        
        # Benefits slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Business Benefits of Autonomous Mode"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        benefits = [
            "⏱️ Dramatically reduced time-to-value (weeks → minutes)",
            "🎯 Optimal performance without manual tuning",
            "🔄 Robust operations with automatic adaptation", 
            "👨‍💼 Expert-level decisions accessible to non-experts",
            "📈 Scalable deployment across any infrastructure",
            "🔮 Future-proof architecture that evolves",
            "💰 Reduced operational costs and expertise requirements",
            "🛡️ Consistent, reliable anomaly detection"
        ]
        
        text_frame.text = benefits[0]
        for benefit in benefits[1:]:
            p = text_frame.add_paragraph()
            p.text = benefit
            p.font.size = Pt(16)
    
    def create_selection_guide_slides(self):
        """Create slides for algorithm selection guide."""
        self.create_section_divider("Algorithm Selection", "Systematic Decision Framework")
        
        # Selection criteria slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Multi-Criteria Selection Framework"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        criteria = [
            "📊 Data Characteristics (30%): Size, dimensionality, type, distribution",
            "🎯 Performance Requirements (25%): Accuracy, precision, recall, F1-score",
            "💻 Computational Constraints (20%): Training time, speed, memory",
            "🔍 Interpretability Needs (15%): Explainability, transparency, trust",
            "🏢 Domain Requirements (10%): Compliance, regulations, standards"
        ]
        
        text_frame.text = criteria[0]
        for criterion in criteria[1:]:
            p = text_frame.add_paragraph()
            p.text = criterion
            p.font.size = Pt(16)
        
        # Decision guidelines slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Expert Decision Guidelines"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        guidelines = [
            "🔍 Data-Driven Decisions: Let data characteristics guide selection",
            "📊 Systematic Evaluation: Use structured frameworks consistently",
            "📈 Progressive Complexity: Start simple, increase when needed",
            "⚖️ Multi-Objective Balance: Accuracy, speed, interpretability, resources",
            "🏢 Domain Integration: Incorporate expertise and constraints",
            "🔄 Continuous Learning: Adapt based on performance feedback"
        ]
        
        text_frame.text = guidelines[0]
        for guideline in guidelines[1:]:
            p = text_frame.add_paragraph()
            p.text = guideline
            p.font.size = Pt(16)
    
    def create_business_procedures_slides(self):
        """Create slides for business procedures section."""
        self.create_section_divider("Business Procedures", "Monthly Data Quality Testing Framework")
        
        # Monthly testing overview slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Monthly Testing Framework"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        framework = [
            "📅 Week 1: Data Collection and Validation",
            "🔍 Week 2: Anomaly Detection and Analysis",
            "📊 Week 3: Trend Analysis and Reporting",
            "📋 Week 4: Review, Documentation, and Planning"
        ]
        
        text_frame.text = framework[0]
        for week in framework[1:]:
            p = text_frame.add_paragraph()
            p.text = week
            p.font.size = Pt(18)
        
        # Quality dimensions slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Data Quality Assessment Dimensions"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        dimensions = [
            "✅ Completeness (30%): Missing values, record counts",
            "🎯 Accuracy (50%): Business rule compliance, validation",
            "⏰ Timeliness (20%): Data freshness, update frequency",
            "🔄 Consistency: Cross-system data alignment",
            "📊 Validity: Format and range compliance",
            "🔗 Integrity: Referential consistency"
        ]
        
        text_frame.text = dimensions[0]
        for dimension in dimensions[1:]:
            p = text_frame.add_paragraph()
            p.text = dimension
            p.font.size = Pt(16)
        
        # KPIs slide
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Key Performance Indicators"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        kpis = [
            "📊 Overall Data Quality Score: >95%",
            "✅ Data Completeness Rate: >98%", 
            "🎯 Data Accuracy Rate: >99%",
            "⏱️ Testing Completion Time: <5 days",
            "🚨 False Positive Rate: <5%",
            "🛠️ Issue Resolution Time: <24 hours",
            "😊 Stakeholder Satisfaction: >4.0/5",
            "✅ Compliance Score: 100%"
        ]
        
        text_frame.text = kpis[0]
        for kpi in kpis[1:]:
            p = text_frame.add_paragraph()
            p.text = kpi
            p.font.size = Pt(16)
    
    def create_implementation_roadmap_slide(self):
        """Create implementation roadmap slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Implementation Roadmap"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        phases = [
            "🚀 Phase 1 (Weeks 1-2): Environment Setup & Data Integration",
            "🔍 Phase 2 (Weeks 3-4): Initial Anomaly Detection & Validation",
            "⚙️ Phase 3 (Weeks 5-6): Algorithm Tuning & Optimization",
            "📊 Phase 4 (Weeks 7-8): Reporting & Dashboard Implementation",
            "🏭 Phase 5 (Weeks 9-10): Production Deployment & Monitoring",
            "🔄 Phase 6 (Ongoing): Continuous Improvement & Maintenance"
        ]
        
        text_frame.text = phases[0]
        for phase in phases[1:]:
            p = text_frame.add_paragraph()
            p.text = phase
            p.font.size = Pt(16)
    
    def create_conclusion_slide(self):
        """Create conclusion slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps & Recommendations"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        next_steps = [
            "🎯 Start with autonomous mode for quick wins",
            "📊 Implement monthly testing procedures",
            "🔍 Focus on high-business-impact data sources first",
            "👥 Train business users on interpretation and escalation",
            "📈 Establish baseline metrics and continuous monitoring",
            "🔄 Plan for iterative improvement based on results",
            "🤝 Engage stakeholders in regular review processes",
            "📚 Leverage knowledge base for continuous learning"
        ]
        
        text_frame.text = next_steps[0]
        for step in next_steps[1:]:
            p = text_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(16)
    
    def create_qa_slide(self):
        """Create Q&A slide."""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.presentation.slide_width,
            self.presentation.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['accent']
        background.line.fill.background()
        
        # Add Q&A text
        qa_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.33), Inches(3.5)
        )
        qa_frame = qa_box.text_frame
        qa_frame.text = "Questions & Discussion"
        qa_para = qa_frame.paragraphs[0]
        qa_para.font.size = Pt(60)
        qa_para.font.color.rgb = self.colors['white']
        qa_para.alignment = PP_ALIGN.CENTER
        
        # Add contact info
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5), Inches(11.33), Inches(1.5)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "For more information: docs.pynomaly.com\nContact: support@pynomaly.com"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = self.colors['white']
        contact_para.alignment = PP_ALIGN.CENTER
        contact_frame.paragraphs[1].font.size = Pt(18)
        contact_frame.paragraphs[1].font.color.rgb = self.colors['white']
        contact_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    def add_footer(self, slide, text: str):
        """Add footer to slide."""
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7), Inches(12.33), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = text
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(10)
        footer_para.font.color.rgb = self.colors['secondary']
        footer_para.alignment = PP_ALIGN.CENTER
    
    def generate_presentation(self, output_path: str):
        """Generate the complete presentation."""
        print("Generating Pynomaly documentation presentation...")
        
        # Create all slides
        self.create_title_slide()
        self.create_agenda_slide()
        self.create_process_overview_slides()
        self.create_architecture_slides()
        self.create_algorithms_slides()
        self.create_autonomous_mode_slides()
        self.create_selection_guide_slides()
        self.create_business_procedures_slides()
        self.create_implementation_roadmap_slide()
        self.create_conclusion_slide()
        self.create_qa_slide()
        
        # Save presentation
        self.presentation.save(output_path)
        print(f"Presentation saved to: {output_path}")
        
        return output_path

def main():
    """Main function to generate the presentation."""
    parser = argparse.ArgumentParser(description="Generate Pynomaly documentation presentation")
    parser.add_argument(
        "--output", 
        default="/mnt/c/Users/andre/Pynomaly/docs/comprehensive/Pynomaly_Documentation_Presentation.pptx",
        help="Output path for the presentation file"
    )
    
    args = parser.parse_args()
    
    # Ensure output directory exists
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Generate presentation
    generator = PynomlalyPresentationGenerator()
    generated_file = generator.generate_presentation(str(output_path))
    
    print(f"✅ Successfully generated presentation: {generated_file}")
    print(f"📊 Total slides: {len(generator.presentation.slides)}")
    
    return generated_file

if __name__ == "__main__":
    main()