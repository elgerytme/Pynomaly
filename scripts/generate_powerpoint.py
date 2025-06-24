#!/usr/bin/env python3
"""
Generate PowerPoint presentation from Banking Anomaly Detection Guide
"""

import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx library not available. Please install it first.")
    print("Run: poetry add python-pptx")
    sys.exit(1)


def create_banking_anomaly_presentation():
    """Create comprehensive PowerPoint presentation for banking anomaly detection."""
    
    # Create presentation
    prs = Presentation()
    
    # Define color scheme (Banking Blue theme)
    primary_blue = RGBColor(0, 51, 102)     # Dark blue
    accent_blue = RGBColor(0, 102, 204)     # Medium blue  
    light_blue = RGBColor(173, 216, 230)    # Light blue
    white = RGBColor(255, 255, 255)
    dark_gray = RGBColor(64, 64, 64)
    
    def add_title_slide(title, subtitle=""):
        """Add a title slide with banking theme."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = primary_blue
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = accent_blue
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
        
        return slide
    
    def add_content_slide(title, content_points):
        """Add a content slide with bullet points."""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = primary_blue
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        tf = content_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = dark_gray
        
        return slide
    
    def add_section_slide(section_title):
        """Add a section divider slide."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = section_title
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.color.rgb = white
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Add background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = primary_blue
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Anomaly and Outlier Detection in Banking",
        "A Business Guide to Modern Risk Detection Technology"
    )
    
    # Slide 2: Executive Summary
    add_content_slide(
        "Executive Summary",
        [
            "Banks process millions of transactions daily across multiple channels",
            "Traditional rule-based systems struggle with volume and sophistication",
            "Anomaly detection provides automated identification of unusual patterns",
            "Pynomaly + PyOD offer enterprise-ready fraud detection solutions",
            "ROI of 300-500% achievable within 3 years of implementation"
        ]
    )
    
    # Section: Understanding Anomaly Detection
    add_section_slide("Understanding Anomaly Detection")
    
    # Slide 3: What is Anomaly Detection?
    add_content_slide(
        "What is Anomaly Detection?",
        [
            "Like a highly trained security guard who knows what 'normal' looks like",
            "Automatically monitors millions of transactions and behaviors",
            "Identifies patterns that deviate from established norms",
            "Works in real-time to prevent fraud before it impacts customers",
            "Adapts and learns from new patterns and analyst feedback"
        ]
    )
    
    # Slide 4: Types of Banking Anomalies
    add_content_slide(
        "Types of Banking Anomalies",
        [
            "Point Anomalies: Single unusual events ($50K withdrawal from $500 account)",
            "Contextual Anomalies: Normal actions in wrong context (3AM rural ATM)",
            "Collective Anomalies: Groups of related events (coordinated attacks)",
            "Behavioral Anomalies: Changes in customer patterns over time",
            "Network Anomalies: Suspicious connections between accounts"
        ]
    )
    
    # Slide 5: Why Banks Need This Technology
    add_content_slide(
        "Why Banks Need Anomaly Detection",
        [
            "Financial Impact: Global banking fraud losses exceed $30B annually",
            "Regulatory Requirements: AML, KYC, fraud prevention mandates",
            "Volume Challenge: Cannot manually review millions of transactions",
            "Sophistication: Fraudsters adapt faster than rule updates",
            "Customer Experience: Reduce false positives that block legitimate transactions"
        ]
    )
    
    # Section: Technology Deep Dive
    add_section_slide("Technology Deep Dive")
    
    # Slide 6: How Anomaly Detection Works
    add_content_slide(
        "How Anomaly Detection Works",
        [
            "Training Phase: Learn patterns from historical transaction data",
            "Pattern Recognition: Understand what constitutes 'normal' behavior",
            "Real-time Scoring: Assign risk scores to incoming transactions",
            "Alert Generation: Flag high-risk transactions for review",
            "Continuous Learning: Adapt based on analyst feedback and outcomes"
        ]
    )
    
    # Slide 7: Pynomaly Platform Overview
    add_content_slide(
        "Pynomaly: Enterprise Anomaly Detection",
        [
            "Multi-Algorithm Integration: Combines PyOD, TODS, PyGOD technologies",
            "Clean Architecture: Modular, scalable, maintainable design",
            "Production-Ready: 99.9% uptime, sub-second response times",
            "Banking-Specific: Risk scoring, regulatory reporting, audit trails",
            "Real-time Processing: Handles millions of transactions per day"
        ]
    )
    
    # Slide 8: PyOD Detection Engine
    add_content_slide(
        "PyOD: The Detection Engine",
        [
            "40+ Advanced Algorithms: Isolation Forest, LOF, One-Class SVM, etc.",
            "Automatic Selection: Chooses optimal algorithms for each data type",
            "Ensemble Methods: Combines multiple algorithms for higher accuracy",
            "Specialized Algorithms: Optimized for different banking use cases",
            "Proven Performance: Widely used in financial services industry"
        ]
    )
    
    # Slide 9: Autonomous Detection Mode
    add_content_slide(
        "Autonomous Detection Mode",
        [
            "Continuous Learning: Adapts to changing patterns automatically",
            "Self-Tuning: Optimizes parameters without manual intervention",
            "Intelligent Prioritization: Routes alerts based on risk and impact",
            "Automated Response: Blocks transactions and initiates workflows",
            "Human Oversight: Maintains controls and escalation procedures"
        ]
    )
    
    # Section: Risk Management
    add_section_slide("Risk Scoring & Prioritization")
    
    # Slide 10: Risk Scoring Framework
    add_content_slide(
        "Risk Scoring Framework",
        [
            "0-30: Low Risk - Routine monitoring and automated handling",
            "31-60: Medium Risk - Scheduled review by standard analysts",
            "61-80: High Risk - Priority investigation by experienced analysts",
            "81-100: Critical Risk - Immediate action and senior analyst review",
            "Multi-factor Scoring: Statistical, contextual, network, and intelligence"
        ]
    )
    
    # Slide 11: Alert Prioritization
    add_content_slide(
        "Intelligent Alert Prioritization",
        [
            "Department-Specific Views: Fraud, AML, Risk, Operations, Compliance",
            "Dynamic Filtering: Geographic, product, channel, and time-based filters",
            "Alert Clustering: Group related suspicious activities together",
            "Predictive Analytics: Anticipate trends before they fully emerge",
            "Feedback Loop: Learn from analyst decisions to improve accuracy"
        ]
    )
    
    # Section: Department Applications
    add_section_slide("Department-Specific Applications")
    
    # Slide 12: Fraud Prevention Department
    add_content_slide(
        "Fraud Prevention Applications",
        [
            "Real-time Transaction Monitoring across all payment channels",
            "Account Takeover Detection through login pattern analysis",
            "Card Fraud Prevention with spending pattern recognition",
            "Digital Banking Fraud for online and mobile anomalies",
            "60-80% reduction in false positives, faster case resolution"
        ]
    )
    
    # Slide 13: AML & Compliance
    add_content_slide(
        "AML & Compliance Applications",
        [
            "Suspicious Activity Reporting (SAR) automation",
            "Money Laundering Pattern Detection (structuring, layering)",
            "Network Analysis for criminal organization identification",
            "Regulatory Compliance automation and audit trail generation",
            "Risk-based Customer Due Diligence and monitoring"
        ]
    )
    
    # Slide 14: Risk Management
    add_content_slide(
        "Risk Management Applications",
        [
            "Credit Risk Early Warning through customer behavior changes",
            "Market Risk Detection for unusual trading patterns",
            "Operational Risk monitoring of systems and processes",
            "Model Risk Management and performance monitoring",
            "Portfolio-wide risk assessment and stress testing"
        ]
    )
    
    # Section: Implementation
    add_section_slide("Implementation Strategy")
    
    # Slide 15: Implementation Roadmap
    add_content_slide(
        "Phased Implementation Approach",
        [
            "Phase 1 (Months 1-3): Foundation - Basic fraud detection",
            "Phase 2 (Months 4-6): Enhancement - Multi-department expansion",
            "Phase 3 (Months 7-12): Autonomous - Full automation deployment",
            "Phase 4 (Ongoing): Improvement - Continuous optimization",
            "Risk-managed approach with clear milestones and success metrics"
        ]
    )
    
    # Slide 16: Success Factors
    add_content_slide(
        "Implementation Success Factors",
        [
            "Executive Sponsorship: Strong C-level support and resource commitment",
            "Cross-Department Collaboration: Involvement of all stakeholders",
            "Change Management: Comprehensive training and communication",
            "Technical Excellence: Proper integration and performance optimization",
            "Continuous Improvement: Ongoing adaptation and enhancement"
        ]
    )
    
    # Section: Business Benefits
    add_section_slide("Business Benefits & ROI")
    
    # Slide 17: Quantifiable Benefits
    add_content_slide(
        "Quantifiable Business Benefits",
        [
            "Fraud Detection: 85-95% accuracy vs. 60-70% traditional systems",
            "Cost Reduction: 70% less investigation time, 60-80% fewer false positives",
            "Compliance: Avoid millions in AML fines, 90% automated reporting",
            "Customer Experience: 80% reduction in legitimate transaction blocks",
            "Operational Efficiency: 40-50% analyst time redirected to high-value work"
        ]
    )
    
    # Slide 18: ROI Analysis
    add_content_slide(
        "Return on Investment",
        [
            "Investment: $3M-$8M over 3 years (varies by bank size)",
            "Annual Benefits: $9M-$25M (fraud reduction, operations, compliance)",
            "Year 1 ROI: 100-200%",
            "Year 2 ROI: 200-400%", 
            "Year 3+ ROI: 300-500%"
        ]
    )
    
    # Slide 19: Strategic Benefits
    add_content_slide(
        "Strategic Business Value",
        [
            "Competitive Advantage: Technology leadership in fraud prevention",
            "Regulatory Positioning: Proactive compliance demonstrates commitment",
            "Business Agility: Rapid adaptation to new threats and growth",
            "Risk Management: Comprehensive enterprise risk visibility",
            "Innovation Platform: Foundation for advanced analytics and AI"
        ]
    )
    
    # Slide 20: Conclusion
    add_content_slide(
        "Key Takeaways",
        [
            "Anomaly detection is transformational for modern banking security",
            "Pynomaly + PyOD provide enterprise-ready, proven technology",
            "ROI of 300-500% achievable with proper implementation",
            "Phased approach reduces risk and ensures successful adoption",
            "Investment today positions bank for competitive advantage tomorrow"
        ]
    )
    
    # Slide 21: Next Steps
    add_content_slide(
        "Recommended Next Steps",
        [
            "Executive briefing and strategic alignment session",
            "Technical architecture review and system integration planning",
            "Pilot program design for initial fraud detection use case",
            "Resource planning and team training requirements assessment",
            "Detailed business case development and budget approval process"
        ]
    )
    
    return prs


def main():
    """Generate the PowerPoint presentation."""
    print("Generating Banking Anomaly Detection PowerPoint presentation...")
    
    try:
        # Create presentation
        prs = create_banking_anomaly_presentation()
        
        # Save presentation
        output_dir = Path(__file__).parent.parent / "docs"
        output_file = output_dir / "Banking_Anomaly_Detection_Presentation.pptx"
        
        # Ensure output directory exists
        output_dir.mkdir(exist_ok=True)
        
        # Save the presentation
        prs.save(str(output_file))
        
        print(f"PowerPoint presentation saved to: {output_file}")
        print(f"Presentation contains {len(prs.slides)} slides")
        
    except Exception as e:
        print(f"Error generating presentation: {e}")
        return 1
    
    return 0


if __name__ == "__main__":
    sys.exit(main())